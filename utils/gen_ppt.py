from pptx import Presentation
from pptx.util import Inches
from common import get_path
## global and constants
EMPTY_LAYOUT_NUM = 10

def gen_ppt(summary):
  prs = open_template()
  gen_summary(prs, summary)
  gen_charts(prs, summary.result_count)
  prs.save(r'./outputs/Report.pptx')

def open_template():
  prs = Presentation(r'./utils/template.pptx')
  return prs

def gen_charts(prs, num):
  for i in range(2, num + 2, 1):
    EMPTY_LAYOUT = prs.slide_layouts[EMPTY_LAYOUT_NUM]
    slide = prs.slides.add_slide(EMPTY_LAYOUT)
    title = slide.shapes.title
    title.text = "PCIE Test " + str(i-1)
    picture = slide.shapes.add_picture(get_path("\\outputs\\img\\" + str(i) +".png"), Inches(2), Inches(2))

def gen_summary(prs, summary):
  EMPTY_LAYOUT = prs.slide_layouts[EMPTY_LAYOUT_NUM]
  slide = prs.slides.add_slide(EMPTY_LAYOUT)
  title = slide.shapes.title
  title.text = "PCIE Test Summary"
  x, y, cx, cy = Inches(4), Inches(1), Inches(2), Inches(1)
  shape = slide.shapes.add_table(2, 2, x, y, cx, cy)
  table = shape.table


  table.cell(0, 0).text = 'Pass'
  table.cell(0, 1).text = 'Fail'
  table.cell(1, 0).text = str(summary.pass_count)
  table.cell(1, 1).text = str(summary.fail_count)

  picture = slide.shapes.add_picture(get_path("\\outputs\\img\\1.png"), Inches(2), Inches(3))
# gen_ppt()